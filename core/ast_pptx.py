"""
core/ast_pptx.py — Parser Recursivo e AST do PowerPoint

Converte um arquivo .pptx em uma Árvore de Sintaxe Abstrata (AST)
hierárquica e serializável em JSON, capturando:
  - Metadados da apresentação
  - Cada slide com seus shapes
  - Cada shape com geometria, formatação e texto
  - Cada parágrafo e run com estilos exatos
  - Tabelas, imagens, grupos e SmartArts

Classes:
  PptxAST          — raiz da árvore
  SlideNode        — nó de um slide
  ShapeNode        — nó de um shape genérico
  TextFrameNode    — nó do frame de texto
  ParagraphNode    — nó de parágrafo
  RunNode          — nó de run de texto
  TableNode        — nó de tabela
  ImageNode        — nó de imagem
  GroupNode        — nó de grupo
"""

from __future__ import annotations

import hashlib
import json
from dataclasses import dataclass, field, asdict
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.shapes.picture import Picture
from pptx.table import Table
from pptx.text.text import TextFrame, _Paragraph, _Run


# ── Enums ────────────────────────────────────────────────────────────────────

class NodeType(str, Enum):
    PRESENTATION = "presentation"
    SLIDE = "slide"
    SHAPE = "shape"
    TEXT_FRAME = "text_frame"
    PARAGRAPH = "paragraph"
    RUN = "run"
    TABLE = "table"
    TABLE_ROW = "table_row"
    TABLE_CELL = "table_cell"
    IMAGE = "image"
    GROUP = "group"
    PLACEHOLDER = "placeholder"


class AlignEnum(str, Enum):
    LEFT = "left"
    CENTER = "center"
    RIGHT = "right"
    JUSTIFY = "justify"
    UNKNOWN = "unknown"


# ── Helpers ──────────────────────────────────────────────────────────────────

def _emu_to_inches(emu: Optional[int]) -> Optional[float]:
    if emu is None:
        return None
    return round(emu / 914400, 4)


def _pt_to_float(pt_val) -> Optional[float]:
    if pt_val is None:
        return None
    try:
        return round(pt_val.pt, 2)
    except Exception:
        return None


def _rgb_to_hex(rgb: Optional[RGBColor]) -> Optional[str]:
    if rgb is None:
        return None
    return f"#{rgb.red:02X}{rgb.green:02X}{rgb.blue:02X}"


def _align_name(align) -> str:
    mapping = {
        PP_ALIGN.LEFT:    AlignEnum.LEFT,
        PP_ALIGN.CENTER:  AlignEnum.CENTER,
        PP_ALIGN.RIGHT:   AlignEnum.RIGHT,
        PP_ALIGN.JUSTIFY: AlignEnum.JUSTIFY,
    }
    return mapping.get(align, AlignEnum.UNKNOWN).value


# ── Nós da AST ───────────────────────────────────────────────────────────────

@dataclass
class RunNode:
    """Nó atômico de texto — um run dentro de um parágrafo."""
    node_type: str = NodeType.RUN.value
    text: str = ""
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    font_name: Optional[str] = None
    font_size_pt: Optional[float] = None
    font_color_hex: Optional[str] = None
    strike: Optional[bool] = None
    highlight_color_hex: Optional[str] = None

    @classmethod
    def from_run(cls, run: _Run) -> "RunNode":
        font = run.font
        color_hex = None
        try:
            if font.color and font.color.type is not None:
                color_hex = _rgb_to_hex(font.color.rgb)
        except Exception:
            pass

        return cls(
            text=run.text,
            bold=font.bold,
            italic=font.italic,
            underline=font.underline,
            font_name=font.name,
            font_size_pt=_pt_to_float(font.size),
            font_color_hex=color_hex,
            strike=font.strike,
        )

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


@dataclass
class ParagraphNode:
    """Nó de parágrafo — contém runs e formatação de parágrafo."""
    node_type: str = NodeType.PARAGRAPH.value
    text: str = ""          # texto completo concatenado
    runs: List[RunNode] = field(default_factory=list)
    alignment: str = AlignEnum.UNKNOWN.value
    level: int = 0          # nível de indentação/bullet
    space_before_pt: Optional[float] = None
    space_after_pt: Optional[float] = None
    line_spacing: Optional[float] = None
    bullet: bool = False
    bullet_char: Optional[str] = None

    @classmethod
    def from_paragraph(cls, para: _Paragraph) -> "ParagraphNode":
        runs = [RunNode.from_run(r) for r in para.runs]
        full_text = "".join(r.text for r in para.runs)

        # Bullet detection
        pPr = para._p.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
        has_bullet = False
        bullet_char = None
        if pPr is not None:
            buChar = pPr.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar")
            buNone = pPr.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone")
            buAutoNum = pPr.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum")
            if buChar is not None:
                has_bullet = True
                bullet_char = buChar.get("char", "•")
            elif buAutoNum is not None:
                has_bullet = True
                bullet_char = "1."

        align_str = AlignEnum.UNKNOWN.value
        try:
            align_str = _align_name(para.alignment)
        except Exception:
            pass

        space_before = None
        space_after = None
        try:
            space_before = _pt_to_float(para.space_before)
            space_after = _pt_to_float(para.space_after)
        except Exception:
            pass

        return cls(
            text=full_text,
            runs=runs,
            alignment=align_str,
            level=para.level,
            space_before_pt=space_before,
            space_after_pt=space_after,
            bullet=has_bullet,
            bullet_char=bullet_char,
        )

    def to_dict(self) -> Dict[str, Any]:
        d = asdict(self)
        d["runs"] = [r.to_dict() for r in self.runs]
        return d


@dataclass
class TextFrameNode:
    """Nó de TextFrame — container de parágrafos dentro de um shape."""
    node_type: str = NodeType.TEXT_FRAME.value
    full_text: str = ""
    paragraphs: List[ParagraphNode] = field(default_factory=list)
    word_wrap: Optional[bool] = None
    auto_size: Optional[str] = None
    margin_left_in: Optional[float] = None
    margin_right_in: Optional[float] = None
    margin_top_in: Optional[float] = None
    margin_bottom_in: Optional[float] = None

    @classmethod
    def from_text_frame(cls, tf: TextFrame) -> "TextFrameNode":
        paragraphs = [ParagraphNode.from_paragraph(p) for p in tf.paragraphs]
        full_text = "\n".join(p.text for p in paragraphs)

        auto_size_name = None
        try:
            auto_size_name = str(tf.auto_size) if tf.auto_size else None
        except Exception:
            pass

        return cls(
            full_text=full_text,
            paragraphs=paragraphs,
            word_wrap=tf.word_wrap,
            auto_size=auto_size_name,
            margin_left_in=_emu_to_inches(tf.margin_left),
            margin_right_in=_emu_to_inches(tf.margin_right),
            margin_top_in=_emu_to_inches(tf.margin_top),
            margin_bottom_in=_emu_to_inches(tf.margin_bottom),
        )

    def to_dict(self) -> Dict[str, Any]:
        d = asdict(self)
        d["paragraphs"] = [p.to_dict() for p in self.paragraphs]
        return d


@dataclass
class ImageNode:
    """Nó de imagem."""
    node_type: str = NodeType.IMAGE.value
    shape_id: int = 0
    name: str = ""
    content_type: str = ""
    width_in: Optional[float] = None
    height_in: Optional[float] = None
    left_in: Optional[float] = None
    top_in: Optional[float] = None
    image_hash: Optional[str] = None

    @classmethod
    def from_picture(cls, pic: Picture) -> "ImageNode":
        img_hash = None
        try:
            img_hash = hashlib.md5(pic.image.blob).hexdigest()
        except Exception:
            pass

        return cls(
            shape_id=pic.shape_id,
            name=pic.name,
            content_type=pic.image.content_type if hasattr(
                pic, "image") else "",
            width_in=_emu_to_inches(pic.width),
            height_in=_emu_to_inches(pic.height),
            left_in=_emu_to_inches(pic.left),
            top_in=_emu_to_inches(pic.top),
            image_hash=img_hash,
        )

    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)


@dataclass
class TableCellNode:
    """Nó de célula de tabela."""
    node_type: str = NodeType.TABLE_CELL.value
    text: str = ""
    row_idx: int = 0
    col_idx: int = 0
    text_frame: Optional[TextFrameNode] = None
    fill_color_hex: Optional[str] = None
    is_merge_origin: bool = False

    def to_dict(self) -> Dict[str, Any]:
        d = asdict(self)
        if self.text_frame:
            d["text_frame"] = self.text_frame.to_dict()
        return d


@dataclass
class TableNode:
    """Nó de tabela."""
    node_type: str = NodeType.TABLE.value
    shape_id: int = 0
    name: str = ""
    rows: int = 0
    cols: int = 0
    width_in: Optional[float] = None
    height_in: Optional[float] = None
    left_in: Optional[float] = None
    top_in: Optional[float] = None
    cells: List[List[TableCellNode]] = field(default_factory=list)

    @classmethod
    def from_table_shape(cls, shape) -> "TableNode":
        table: Table = shape.table
        rows = len(table.rows)
        cols = len(table.columns)
        cells = []
        for r_idx, row in enumerate(table.rows):
            row_cells = []
            for c_idx, cell in enumerate(row.cells):
                tf_node = None
                try:
                    tf_node = TextFrameNode.from_text_frame(cell.text_frame)
                except Exception:
                    pass

                fill_hex = None
                try:
                    fill_hex = _rgb_to_hex(cell.fill.fore_color.rgb)
                except Exception:
                    pass

                row_cells.append(TableCellNode(
                    text=cell.text,
                    row_idx=r_idx,
                    col_idx=c_idx,
                    text_frame=tf_node,
                    fill_color_hex=fill_hex,
                ))
            cells.append(row_cells)

        return cls(
            shape_id=shape.shape_id,
            name=shape.name,
            rows=rows,
            cols=cols,
            width_in=_emu_to_inches(shape.width),
            height_in=_emu_to_inches(shape.height),
            left_in=_emu_to_inches(shape.left),
            top_in=_emu_to_inches(shape.top),
            cells=cells,
        )

    def to_dict(self) -> Dict[str, Any]:
        d = {
            "node_type": self.node_type,
            "shape_id": self.shape_id,
            "name": self.name,
            "rows": self.rows,
            "cols": self.cols,
            "width_in": self.width_in,
            "height_in": self.height_in,
            "left_in": self.left_in,
            "top_in": self.top_in,
            "cells": [[c.to_dict() for c in row] for row in self.cells],
        }
        return d


@dataclass
class ShapeNode:
    """
    Nó genérico de shape — representa qualquer forma do slide.
    Suporta: textbox, placeholder, autoshape, connector.
    """
    node_type: str = NodeType.SHAPE.value
    shape_id: int = 0
    name: str = ""
    shape_type: str = ""
    is_placeholder: bool = False
    placeholder_type: Optional[str] = None
    placeholder_idx: Optional[int] = None
    width_in: Optional[float] = None
    height_in: Optional[float] = None
    left_in: Optional[float] = None
    top_in: Optional[float] = None
    rotation: float = 0.0
    fill_color_hex: Optional[str] = None
    line_color_hex: Optional[str] = None
    line_width_pt: Optional[float] = None
    visible: bool = True
    text_frame: Optional[TextFrameNode] = None

    @classmethod
    def from_shape(cls, shape) -> "ShapeNode":
        tf_node = None
        if shape.has_text_frame:
            try:
                tf_node = TextFrameNode.from_text_frame(shape.text_frame)
            except Exception:
                pass

        is_ph = shape.is_placeholder
        ph_type = None
        ph_idx = None
        if is_ph:
            try:
                ph = shape.placeholder_format
                ph_type = str(ph.type)
                ph_idx = ph.idx
            except Exception:
                pass

        fill_hex = None
        try:
            fill = shape.fill
            if fill.type is not None:
                fill_hex = _rgb_to_hex(fill.fore_color.rgb)
        except Exception:
            pass

        line_hex = None
        line_width = None
        try:
            line = shape.line
            if line.color and line.color.type is not None:
                line_hex = _rgb_to_hex(line.color.rgb)
            line_width = _pt_to_float(line.width) if line.width else None
        except Exception:
            pass

        rotation = 0.0
        try:
            rotation = shape.rotation or 0.0
        except Exception:
            pass

        return cls(
            shape_id=shape.shape_id,
            name=shape.name,
            shape_type=str(shape.shape_type),
            is_placeholder=is_ph,
            placeholder_type=ph_type,
            placeholder_idx=ph_idx,
            width_in=_emu_to_inches(shape.width),
            height_in=_emu_to_inches(shape.height),
            left_in=_emu_to_inches(shape.left),
            top_in=_emu_to_inches(shape.top),
            rotation=rotation,
            fill_color_hex=fill_hex,
            line_color_hex=line_hex,
            line_width_pt=line_width,
            text_frame=tf_node,
        )

    def to_dict(self) -> Dict[str, Any]:
        d = asdict(self)
        if self.text_frame:
            d["text_frame"] = self.text_frame.to_dict()
        return d


@dataclass
class GroupNode:
    """Nó de grupo de shapes."""
    node_type: str = NodeType.GROUP.value
    shape_id: int = 0
    name: str = ""
    width_in: Optional[float] = None
    height_in: Optional[float] = None
    left_in: Optional[float] = None
    top_in: Optional[float] = None
    # ShapeNode | ImageNode | GroupNode
    children: List[Any] = field(default_factory=list)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "node_type": self.node_type,
            "shape_id": self.shape_id,
            "name": self.name,
            "width_in": self.width_in,
            "height_in": self.height_in,
            "left_in": self.left_in,
            "top_in": self.top_in,
            "children": [c.to_dict() for c in self.children],
        }


@dataclass
class SlideNode:
    """Nó de um slide da apresentação."""
    node_type: str = NodeType.SLIDE.value
    slide_number: int = 0
    slide_id: Optional[int] = None
    layout_name: str = ""
    title: str = ""
    shapes: List[Any] = field(default_factory=list)
    notes: str = ""
    background_color_hex: Optional[str] = None
    has_images: bool = False
    has_tables: bool = False
    has_groups: bool = False
    editable_text_count: int = 0

    def get_text_shapes(self) -> List[ShapeNode]:
        """Retorna apenas shapes com texto."""
        result = []
        for s in self.shapes:
            if isinstance(s, ShapeNode) and s.text_frame:
                result.append(s)
        return result

    def get_shape_by_name(self, name: str) -> Optional[ShapeNode]:
        for s in self.shapes:
            if hasattr(s, "name") and s.name == name:
                return s
        return None

    def get_placeholder_by_idx(self, idx: int) -> Optional[ShapeNode]:
        for s in self.shapes:
            if isinstance(s, ShapeNode) and s.placeholder_idx == idx:
                return s
        return None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "node_type": self.node_type,
            "slide_number": self.slide_number,
            "slide_id": self.slide_id,
            "layout_name": self.layout_name,
            "title": self.title,
            "shapes": [s.to_dict() for s in self.shapes],
            "notes": self.notes,
            "background_color_hex": self.background_color_hex,
            "has_images": self.has_images,
            "has_tables": self.has_tables,
            "has_groups": self.has_groups,
            "editable_text_count": self.editable_text_count,
        }


@dataclass
class PptxAST:
    """
    Raiz da Árvore de Sintaxe Abstrata de uma apresentação .pptx.

    Uso:
        ast = PptxAST.from_file("template.pptx")
        ast.save_json("cache/template_ast.json")

        loaded = PptxAST.load_json("cache/template_ast.json")
    """
    node_type: str = NodeType.PRESENTATION.value
    file_path: str = ""
    file_hash: str = ""
    slide_width_in: float = 13.33
    slide_height_in: float = 7.5
    slide_count: int = 0
    slides: List[SlideNode] = field(default_factory=list)
    core_properties: Dict[str, str] = field(default_factory=dict)
    ast_version: str = "1.0.0"

    # ── Factory ──────────────────────────────────────────────────────────────

    @classmethod
    def from_file(cls, file_path: Union[str, Path]) -> "PptxAST":
        """Parse recursivo de um arquivo .pptx → AST completa."""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"Arquivo não encontrado: {path}")

        # Hash do arquivo para cache
        file_hash = hashlib.md5(path.read_bytes()).hexdigest()

        prs = Presentation(str(path))

        # Core properties
        core_props = {}
        try:
            cp = prs.core_properties
            core_props = {
                "title": cp.title or "",
                "author": cp.author or "",
                "subject": cp.subject or "",
                "keywords": cp.keywords or "",
                "created": str(cp.created) if cp.created else "",
                "modified": str(cp.modified) if cp.modified else "",
            }
        except Exception:
            pass

        slides = []
        for i, slide in enumerate(prs.slides):
            slide_node = cls._parse_slide(slide, slide_number=i + 1)
            slides.append(slide_node)

        return cls(
            file_path=str(path.resolve()),
            file_hash=file_hash,
            slide_width_in=round(prs.slide_width.inches, 4),
            slide_height_in=round(prs.slide_height.inches, 4),
            slide_count=len(slides),
            slides=slides,
            core_properties=core_props,
        )

    @classmethod
    def _parse_slide(cls, slide, slide_number: int) -> SlideNode:
        """Parse recursivo de um slide."""
        layout_name = ""
        try:
            layout_name = slide.slide_layout.name
        except Exception:
            pass

        # Notes
        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
        except Exception:
            pass

        # Parse shapes recursivamente
        parsed_shapes = []
        has_images = False
        has_tables = False
        has_groups = False
        editable_count = 0

        for shape in slide.shapes:
            node = cls._parse_shape(shape)
            if node:
                parsed_shapes.append(node)
                if isinstance(node, ImageNode):
                    has_images = True
                elif isinstance(node, TableNode):
                    has_tables = True
                elif isinstance(node, GroupNode):
                    has_groups = True
                elif isinstance(node, ShapeNode) and node.text_frame:
                    editable_count += 1

        # Título do slide
        title = ""
        for s in parsed_shapes:
            if isinstance(s, ShapeNode) and s.is_placeholder and s.placeholder_idx == 0:
                title = s.text_frame.full_text.strip() if s.text_frame else ""
                break
        if not title:
            for s in parsed_shapes:
                if isinstance(s, ShapeNode) and s.text_frame:
                    t = s.text_frame.full_text.strip()
                    if t:
                        title = t[:80]
                        break

        return SlideNode(
            slide_number=slide_number,
            layout_name=layout_name,
            title=title,
            shapes=parsed_shapes,
            notes=notes_text,
            has_images=has_images,
            has_tables=has_tables,
            has_groups=has_groups,
            editable_text_count=editable_count,
        )

    @classmethod
    def _parse_shape(cls, shape) -> Optional[Union[ShapeNode, ImageNode, TableNode, GroupNode]]:
        """Parser recursivo de shapes — despacha para o tipo correto."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        shape_type = shape.shape_type

        # Imagem
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                return ImageNode.from_picture(shape)
            except Exception:
                return None

        # Tabela
        if shape.has_table:
            try:
                return TableNode.from_table_shape(shape)
            except Exception:
                return None

        # Grupo (recursivo)
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            children = []
            try:
                for child in shape.shapes:
                    child_node = cls._parse_shape(child)
                    if child_node:
                        children.append(child_node)
            except Exception:
                pass

            return GroupNode(
                shape_id=shape.shape_id,
                name=shape.name,
                width_in=_emu_to_inches(shape.width),
                height_in=_emu_to_inches(shape.height),
                left_in=_emu_to_inches(shape.left),
                top_in=_emu_to_inches(shape.top),
                children=children,
            )

        # Shape genérico com ou sem texto
        try:
            return ShapeNode.from_shape(shape)
        except Exception:
            return None

    # ── Serialização ──────────────────────────────────────────────────────────

    def to_dict(self) -> Dict[str, Any]:
        return {
            "node_type": self.node_type,
            "file_path": self.file_path,
            "file_hash": self.file_hash,
            "slide_width_in": self.slide_width_in,
            "slide_height_in": self.slide_height_in,
            "slide_count": self.slide_count,
            "ast_version": self.ast_version,
            "core_properties": self.core_properties,
            "slides": [s.to_dict() for s in self.slides],
        }

    def to_json(self, indent: int = 2) -> str:
        return json.dumps(self.to_dict(), ensure_ascii=False, indent=indent)

    def save_json(self, output_path: Union[str, Path]) -> Path:
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(self.to_json(), encoding="utf-8")
        return path

    @classmethod
    def load_json(cls, json_path: Union[str, Path]) -> "PptxAST":
        """Carrega AST de um JSON serializado."""
        path = Path(json_path)
        data = json.loads(path.read_text(encoding="utf-8"))
        return cls._from_dict(data)

    @classmethod
    def _from_dict(cls, data: Dict[str, Any]) -> "PptxAST":
        """Deserializa AST de um dicionário."""
        slides = [cls._slide_from_dict(s) for s in data.get("slides", [])]
        return cls(
            file_path=data.get("file_path", ""),
            file_hash=data.get("file_hash", ""),
            slide_width_in=data.get("slide_width_in", 13.33),
            slide_height_in=data.get("slide_height_in", 7.5),
            slide_count=data.get("slide_count", 0),
            slides=slides,
            core_properties=data.get("core_properties", {}),
            ast_version=data.get("ast_version", "1.0.0"),
        )

    @classmethod
    def _slide_from_dict(cls, data: Dict[str, Any]) -> SlideNode:
        shapes_raw = data.get("shapes", [])
        shapes = []
        for s in shapes_raw:
            node = cls._shape_from_dict(s)
            if node:
                shapes.append(node)

        return SlideNode(
            slide_number=data.get("slide_number", 0),
            slide_id=data.get("slide_id"),
            layout_name=data.get("layout_name", ""),
            title=data.get("title", ""),
            shapes=shapes,
            notes=data.get("notes", ""),
            background_color_hex=data.get("background_color_hex"),
            has_images=data.get("has_images", False),
            has_tables=data.get("has_tables", False),
            has_groups=data.get("has_groups", False),
            editable_text_count=data.get("editable_text_count", 0),
        )

    @classmethod
    def _shape_from_dict(cls, data: Dict[str, Any]) -> Optional[Any]:
        nt = data.get("node_type", "")
        if nt == NodeType.IMAGE.value:
            return ImageNode(**{k: v for k, v in data.items() if k != "node_type"})
        if nt == NodeType.TABLE.value:
            return TableNode(
                shape_id=data.get("shape_id", 0),
                name=data.get("name", ""),
                rows=data.get("rows", 0),
                cols=data.get("cols", 0),
                width_in=data.get("width_in"),
                height_in=data.get("height_in"),
                left_in=data.get("left_in"),
                top_in=data.get("top_in"),
            )
        if nt == NodeType.GROUP.value:
            children = [cls._shape_from_dict(c)
                        for c in data.get("children", [])]
            return GroupNode(
                shape_id=data.get("shape_id", 0),
                name=data.get("name", ""),
                width_in=data.get("width_in"),
                height_in=data.get("height_in"),
                left_in=data.get("left_in"),
                top_in=data.get("top_in"),
                children=[c for c in children if c],
            )
        # Shape genérico
        tf_data = data.pop("text_frame", None)
        tf_node = None
        if tf_data:
            tf_node = cls._text_frame_from_dict(tf_data)
        shape = ShapeNode(**{k: v for k, v in data.items()
                          if k not in ("node_type", "text_frame")})
        shape.text_frame = tf_node
        return shape

    @classmethod
    def _text_frame_from_dict(cls, data: Dict) -> TextFrameNode:
        paras = [cls._paragraph_from_dict(p)
                 for p in data.get("paragraphs", [])]
        return TextFrameNode(
            full_text=data.get("full_text", ""),
            paragraphs=paras,
            word_wrap=data.get("word_wrap"),
            auto_size=data.get("auto_size"),
            margin_left_in=data.get("margin_left_in"),
            margin_right_in=data.get("margin_right_in"),
            margin_top_in=data.get("margin_top_in"),
            margin_bottom_in=data.get("margin_bottom_in"),
        )

    @classmethod
    def _paragraph_from_dict(cls, data: Dict) -> ParagraphNode:
        runs = [RunNode(**{k: v for k, v in r.items() if k != "node_type"})
                for r in data.get("runs", [])]
        return ParagraphNode(
            text=data.get("text", ""),
            runs=runs,
            alignment=data.get("alignment", AlignEnum.UNKNOWN.value),
            level=data.get("level", 0),
            space_before_pt=data.get("space_before_pt"),
            space_after_pt=data.get("space_after_pt"),
            bullet=data.get("bullet", False),
            bullet_char=data.get("bullet_char"),
        )

    # ── Utilitários ───────────────────────────────────────────────────────────

    def get_slide(self, number: int) -> Optional[SlideNode]:
        """Slide pelo número (1-indexed)."""
        for s in self.slides:
            if s.slide_number == number:
                return s
        return None

    def summary(self) -> str:
        """Resumo textual da AST para debug."""
        lines = [
            f"📊 PPTX AST — {Path(self.file_path).name}",
            f"   Hash: {self.file_hash}",
            f"   Slides: {self.slide_count}",
            f"   Tamanho: {self.slide_width_in}\" × {self.slide_height_in}\"",
            "",
        ]
        for slide in self.slides:
            shapes_info = f"{slide.editable_text_count} text | {'🖼' if slide.has_images else ''} {'📋' if slide.has_tables else ''} {'👥' if slide.has_groups else ''}"
            lines.append(
                f"  Slide {slide.slide_number:2d} [{slide.layout_name:30s}] {shapes_info}")
            if slide.title:
                lines.append(f"           → {slide.title[:70]}")
        return "\n".join(lines)


# ── CLI standalone ───────────────────────────────────────────────────────────
if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Uso: python ast_pptx.py <arquivo.pptx> [output.json]")
        sys.exit(1)

    pptx_file = sys.argv[1]
    json_out = sys.argv[2] if len(sys.argv) > 2 else None

    print(f"Parsing {pptx_file}...")
    ast = PptxAST.from_file(pptx_file)
    print(ast.summary())

    if json_out:
        saved = ast.save_json(json_out)
        print(f"\n✅ AST salva em: {saved}")
