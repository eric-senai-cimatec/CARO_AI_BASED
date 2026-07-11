"""
core/layout_preserver.py — Preserva layouts de slides durante edições PPTX.

Esta classe captura as posições e dimensões dos shapes de cada slide
antes de alterações e permite restaurar o layout original posteriormente.
"""

from __future__ import annotations
from dataclasses import dataclass
from typing import Dict, List, Optional

from pptx.slide import Slide
from pptx.shapes.base import BaseShape


@dataclass
class ShapeLayout:
    shape_id: int
    name: str
    left: int
    top: int
    width: int
    height: int

    def apply_to(self, shape: BaseShape) -> None:
        shape.left = self.left
        shape.top = self.top
        shape.width = self.width
        shape.height = self.height


class LayoutPreserver:
    """Captura e restaura o layout de slides do PPTX."""

    def __init__(self) -> None:
        self._slide_layouts: Dict[int, List[ShapeLayout]] = {}

    def capture_presentation(self, prs) -> None:
        """Captura o layout de todos os slides da apresentação."""
        for index, slide in enumerate(prs.slides, start=1):
            self.capture_slide(slide, slide_number=index)

    def capture_slide(self, slide: Slide, slide_number: Optional[int] = None) -> None:
        """Captura o layout de um slide específico."""
        idx = slide_number if slide_number is not None else self._get_slide_number(slide)
        self._slide_layouts[idx] = [self._shape_layout(shape) for shape in slide.shapes]

    def restore_presentation(self, prs) -> None:
        """Restaura o layout capturado para todos os slides da apresentação."""
        for index, slide in enumerate(prs.slides, start=1):
            self.restore_slide(slide, slide_number=index)

    def restore_slide(self, slide: Slide, slide_number: Optional[int] = None) -> None:
        """Restaura o layout de um slide específico."""
        idx = slide_number if slide_number is not None else self._get_slide_number(slide)
        saved_layout = self._slide_layouts.get(idx)
        if not saved_layout:
            return

        for shape in slide.shapes:
            layout = self._find_layout_for_shape(shape, saved_layout)
            if layout:
                layout.apply_to(shape)

    def clear(self) -> None:
        """Limpa o cache de layout armazenado."""
        self._slide_layouts.clear()

    def _shape_layout(self, shape: BaseShape) -> ShapeLayout:
        return ShapeLayout(
            shape_id=shape.shape_id,
            name=shape.name or "",
            left=shape.left,
            top=shape.top,
            width=shape.width,
            height=shape.height,
        )

    def _find_layout_for_shape(
        self, shape: BaseShape, layouts: List[ShapeLayout]
    ) -> Optional[ShapeLayout]:
        for layout in layouts:
            if shape.shape_id == layout.shape_id or (
                layout.name and shape.name == layout.name
            ):
                return layout
        return None

    def _get_slide_number(self, slide: Slide) -> int:
        try:
            partname = slide.partname
            return int(partname.split("/")[-1].replace("slide", "").replace(".xml", ""))
        except Exception:
            return -1
