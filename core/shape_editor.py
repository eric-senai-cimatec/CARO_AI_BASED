"""
core/shape_editor.py — Editor de Shapes Preservando Formatação

Edita o conteúdo de texto em shapes do .pptx sem destruir:
  - Cores de fonte
  - Tamanhos de fonte
  - Bold / Italic / Underline
  - Alinhamento de parágrafos
  - Bullets e numeração
  - Espaçamento entre parágrafos
  - Formatação de células de tabela

Classes:
  ShapeEditor  — editor principal, opera sobre python-pptx Presentation
  TextReplacer — substitui texto preservando runs de estilo
  TableEditor  — edita células de tabelas
  FormatPreserver — copia formatação de um run para outro
"""

from __future__ import annotations

import copy
import logging
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.shapes.autoshape import Shape
from pptx.text.text import TextFrame, _Paragraph, _Run
from pptx.util import Inches, Pt, Emu
from lxml import etree

logger = logging.getLogger(__name__)


# ── Helpers de formatação ────────────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> Optional[RGBColor]:
    """Converte string hex (#RRGGBB ou RRGGBB) para RGBColor."""
    if not hex_color:
        return None
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return None
    try:
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        return RGBColor(r, g, b)
    except ValueError:
        return None


def _copy_run_format(source_run: _Run, target_run: _Run):
    """
    Copia a formatação de um run para outro:
    cor, tamanho, bold, italic, underline, fonte.
    """
    src_font = source_run.font
    tgt_font = target_run.font

    if src_font.bold is not None:
        tgt_font.bold = src_font.bold
    if src_font.italic is not None:
        tgt_font.italic = src_font.italic
    if src_font.underline is not None:
        tgt_font.underline = src_font.underline
    if src_font.size is not None:
        tgt_font.size = src_font.size
    if src_font.name:
        tgt_font.name = src_font.name
    try:
        if src_font.color and src_font.color.type is not None:
            tgt_font.color.rgb = src_font.color.rgb
    except Exception:
        pass


def _copy_paragraph_format(source_para: _Paragraph, target_para: _Paragraph):
    """Copia formatação de parágrafo: alinhamento, espaçamento, level."""
    try:
        if source_para.alignment:
            target_para.alignment = source_para.alignment
    except Exception:
        pass
    try:
        target_para.level = source_para.level
    except Exception:
        pass
    try:
        if source_para.space_before:
            target_para.space_before = source_para.space_before
        if source_para.space_after:
            target_para.space_after = source_para.space_after
    except Exception:
        pass


# ── FormatPreserver ──────────────────────────────────────────────────────────

class FormatPreserver:
    """
    Captura e restaura a formatação de um TextFrame.
    Usado para salvar o estilo antes de limpar o conteúdo.
    """

    @dataclass
    class RunSnapshot:
        bold: Optional[bool]
        italic: Optional[bool]
        underline: Optional[bool]
        font_size: Optional[Pt]
        font_name: Optional[str]
        color_rgb: Optional[RGBColor]

    @dataclass
    class ParaSnapshot:
        alignment: Any
        level: int
        space_before: Any
        space_after: Any
        runs: List["FormatPreserver.RunSnapshot"]

    def __init__(self, text_frame: TextFrame):
        self.snapshots: List[FormatPreserver.ParaSnapshot] = []
        self._capture(text_frame)

    def _capture(self, tf: TextFrame):
        for para in tf.paragraphs:
            run_snaps = []
            for run in para.runs:
                font = run.font
                color_rgb = None
                try:
                    if font.color and font.color.type is not None:
                        color_rgb = font.color.rgb
                except Exception:
                    pass

                run_snaps.append(self.RunSnapshot(
                    bold=font.bold,
                    italic=font.italic,
                    underline=font.underline,
                    font_size=font.size,
                    font_name=font.name,
                    color_rgb=color_rgb,
                ))

            snap = self.ParaSnapshot(
                alignment=para.alignment,
                level=para.level,
                space_before=para.space_before,
                space_after=para.space_after,
                runs=run_snaps,
            )
            self.snapshots.append(snap)

    def first_run_snapshot(self) -> Optional["FormatPreserver.RunSnapshot"]:
        """Retorna o snapshot do primeiro run disponível."""
        for snap in self.snapshots:
            if snap.runs:
                return snap.runs[0]
        return None

    def first_para_snapshot(self) -> Optional["FormatPreserver.ParaSnapshot"]:
        if self.snapshots:
            return self.snapshots[0]
        return None


# ── TextReplacer ─────────────────────────────────────────────────────────────

class TextReplacer:
    """
    Substitui o texto de um TextFrame preservando o estilo.

    Estratégia:
      1. Captura FormatPreserver do frame original
      2. Limpa todos os parágrafos/runs
      3. Escreve o novo conteúdo usando o estilo do run[0] original
      4. Para listas de bullets, replica a formatação de bullet
    """

    def __init__(self, text_frame: TextFrame):
        self.tf = text_frame
        self.preserver = FormatPreserver(text_frame)

    def replace_with_text(self, new_text: str):
        """
        Substitui conteúdo por texto simples (sem bullets).
        Preserva formatação do primeiro run.
        """
        if not new_text:
            return

        self._clear_frame()
        first_run_snap = self.preserver.first_run_snapshot()
        first_para_snap = self.preserver.first_para_snapshot()

        # Escrever linha a linha
        lines = new_text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                para = self.tf.paragraphs[0]
            else:
                para = self.tf.add_paragraph()

            # Aplicar formatação de parágrafo
            if first_para_snap:
                try:
                    if first_para_snap.alignment:
                        para.alignment = first_para_snap.alignment
                except Exception:
                    pass

            # Criar run com o texto
            run = para.add_run()
            run.text = line

            # Aplicar formatação do run original
            if first_run_snap:
                self._apply_run_snapshot(run, first_run_snap)

    def replace_with_bullets(self, bullets: List[str], bullet_char: str = "•"):
        """
        Substitui conteúdo por uma lista de bullets.
        Preserva formatação do primeiro bullet original.
        """
        if not bullets:
            return

        self._clear_frame()
        first_run_snap = self.preserver.first_run_snapshot()

        for i, text in enumerate(bullets):
            if i == 0:
                para = self.tf.paragraphs[0]
            else:
                para = self.tf.add_paragraph()

            # Configurar bullet XML
            pPr = para._p.get_or_add_pPr()
            # Remover buNone se existir
            for bu_none in pPr.findall(qn("a:buNone")):
                pPr.remove(bu_none)
            # Adicionar buChar
            buChar = etree.SubElement(pPr, qn("a:buChar"))
            buChar.set("char", bullet_char)

            run = para.add_run()
            run.text = text
            if first_run_snap:
                self._apply_run_snapshot(run, first_run_snap)

    def replace_with_key_value(self, pairs: List[Tuple[str, str]], separator: str = ": "):
        """
        Substitui por pares chave:valor — ex: "Empresa: CARO Soluções".
        Chave em bold, valor normal.
        """
        if not pairs:
            return

        self._clear_frame()
        first_run_snap = self.preserver.first_run_snapshot()

        for i, (key, value) in enumerate(pairs):
            if i == 0:
                para = self.tf.paragraphs[0]
            else:
                para = self.tf.add_paragraph()

            # Run chave (bold)
            run_key = para.add_run()
            run_key.text = f"{key}{separator}"
            if first_run_snap:
                self._apply_run_snapshot(run_key, first_run_snap)
            run_key.font.bold = True

            # Run valor (normal)
            run_val = para.add_run()
            run_val.text = value
            if first_run_snap:
                self._apply_run_snapshot(run_val, first_run_snap)
            run_val.font.bold = False

    def _clear_frame(self):
        """Limpa todos os parágrafos do TextFrame, mantendo pelo menos um."""
        # Manter exatamente um parágrafo vazio
        paragraphs = self.tf.paragraphs
        while len(paragraphs) > 1:
            p = paragraphs[-1]._p
            p.getparent().remove(p)

        # Limpar runs do primeiro parágrafo
        first_para = self.tf.paragraphs[0]
        for run in list(first_para.runs):
            run._r.getparent().remove(run._r)

    def _apply_run_snapshot(self, run: _Run, snap: FormatPreserver.RunSnapshot):
        """Aplica um snapshot de formatação a um run."""
        font = run.font
        if snap.bold is not None:
            font.bold = snap.bold
        if snap.italic is not None:
            font.italic = snap.italic
        if snap.underline is not None:
            font.underline = snap.underline
        if snap.font_size is not None:
            font.size = snap.font_size
        if snap.font_name:
            font.name = snap.font_name
        if snap.color_rgb is not None:
            try:
                font.color.rgb = snap.color_rgb
            except Exception:
                pass


# ── TableEditor ───────────────────────────────────────────────────────────────

class TableEditor:
    """Edita células de uma tabela no .pptx."""

    def __init__(self, shape):
        self.shape = shape
        self.table = shape.table

    def set_cell(self, row: int, col: int, text: str, bold: bool = False):
        """Define o texto de uma célula."""
        cell = self.table.cell(row, col)
        tf = cell.text_frame
        replacer = TextReplacer(tf)
        replacer.replace_with_text(text)
        if bold and tf.paragraphs:
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.bold = True

    def fill_row(self, row: int, values: List[str]):
        """Preenche uma linha inteira de colunas."""
        for col, val in enumerate(values):
            if col < len(self.table.columns):
                self.set_cell(row, col, val)

    def fill_all(self, data: List[List[str]]):
        """Preenche toda a tabela com dados 2D."""
        for row_idx, row_data in enumerate(data):
            if row_idx < len(self.table.rows):
                self.fill_row(row_idx, row_data)


# ── ShapeEditor ───────────────────────────────────────────────────────────────

class ShapeEditor:
    """
    Editor principal de shapes do .pptx.

    Opera diretamente sobre o objeto Presentation do python-pptx.
    Preserva toda a formatação original.

    Uso:
        prs = Presentation("template.pptx")
        editor = ShapeEditor(prs)

        # Editar título do slide 4
        editor.set_slide_title(slide_number=4, text="Novo Título")

        # Editar corpo
        editor.set_slide_body(slide_number=4, bullets=[
            "Empresa: CARO Soluções",
            "Contato: João Silva",
        ])

        prs.save("output.pptx")
    """

    def __init__(self, presentation: Presentation):
        self.prs = presentation
        self._slides = list(presentation.slides)

    # ── API de alto nível ─────────────────────────────────────────────────────

    def set_slide_title(self, slide_number: int, text: str):
        """Define o título de um slide (1-indexed)."""
        slide = self._get_slide(slide_number)
        if not slide:
            logger.warning(f"Slide {slide_number} não encontrado")
            return False

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                replacer = TextReplacer(shape.text_frame)
                replacer.replace_with_text(text)
                logger.debug(f"Slide {slide_number} título definido: {text[:50]!r}")
                return True

        # Fallback: primeiro shape de texto grande
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            name_lower = shape.name.lower()
            if "título" in name_lower or "title" in name_lower:
                replacer = TextReplacer(shape.text_frame)
                replacer.replace_with_text(text)
                return True

        logger.warning(f"Shape de título não encontrado no slide {slide_number}")
        return False

    def set_slide_body(
        self,
        slide_number: int,
        text: Optional[str] = None,
        bullets: Optional[List[str]] = None,
        key_values: Optional[List[Tuple[str, str]]] = None,
        shape_name: Optional[str] = None,
    ) -> bool:
        """
        Define o corpo de um slide.

        Prioridade:
          - shape_name: edita shape específico pelo nome
          - bullets: lista de bullets
          - key_values: pares chave:valor
          - text: texto simples
        """
        slide = self._get_slide(slide_number)
        if not slide:
            return False

        # Encontrar shape alvo
        target_shape = None
        if shape_name:
            target_shape = self._find_shape_by_name(slide, shape_name)
        else:
            target_shape = self._find_body_shape(slide)

        if not target_shape or not target_shape.has_text_frame:
            logger.warning(f"Slide {slide_number}: shape de corpo não encontrado")
            return False

        replacer = TextReplacer(target_shape.text_frame)

        if key_values:
            replacer.replace_with_key_value(key_values)
        elif bullets:
            replacer.replace_with_bullets(bullets)
        elif text:
            replacer.replace_with_text(text)

        logger.debug(f"Slide {slide_number} corpo atualizado ({shape_name or 'auto'})")
        return True

    def set_date(self, slide_number: int, date_text: str) -> bool:
        """Define o campo de data de um slide."""
        slide = self._get_slide(slide_number)
        if not slide:
            return False

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            name_lower = shape.name.lower()
            if "data" in name_lower or "date" in name_lower:
                replacer = TextReplacer(shape.text_frame)
                replacer.replace_with_text(date_text)
                return True

        # Tentar placeholder idx 10+ (data)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.idx > 10:
                        replacer = TextReplacer(shape.text_frame)
                        replacer.replace_with_text(date_text)
                        return True
                except Exception:
                    pass

        return False

    def set_shape_text(
        self,
        slide_number: int,
        shape_name: str,
        text: str,
    ) -> bool:
        """Define texto de um shape específico pelo nome."""
        slide = self._get_slide(slide_number)
        if not slide:
            return False
        shape = self._find_shape_by_name(slide, shape_name)
        if not shape or not shape.has_text_frame:
            return False
        TextReplacer(shape.text_frame).replace_with_text(text)
        return True

    def edit_table(self, slide_number: int, data: List[List[str]]) -> bool:
        """Edita a primeira tabela de um slide com dados 2D."""
        slide = self._get_slide(slide_number)
        if not slide:
            return False
        for shape in slide.shapes:
            if shape.has_table:
                TableEditor(shape).fill_all(data)
                return True
        return False

    # ── Operações de slides ───────────────────────────────────────────────────

    def duplicate_slide(self, slide_number: int) -> int:
        """Duplica um slide e retorna o número do novo slide."""
        from pptx.util import Inches
        from lxml import etree

        src_slide = self._get_slide(slide_number)
        if not src_slide:
            raise ValueError(f"Slide {slide_number} não encontrado")

        # Clonar XML do slide
        src_xml = copy.deepcopy(src_slide._element)
        slide_layout = src_slide.slide_layout

        # Criar novo slide com mesmo layout
        new_slide = self.prs.slides.add_slide(slide_layout)
        new_slide._element.getparent().replace(new_slide._element, src_xml)

        new_number = len(self.prs.slides)
        self._slides = list(self.prs.slides)
        return new_number

    def get_slide_text_summary(self, slide_number: int) -> Dict[str, str]:
        """Retorna dicionário shape_name → texto atual do slide."""
        slide = self._get_slide(slide_number)
        if not slide:
            return {}
        result = {}
        for shape in slide.shapes:
            if shape.has_text_frame:
                result[shape.name] = shape.text_frame.text.strip()
        return result

    def list_editable_shapes(self, slide_number: int) -> List[Dict[str, Any]]:
        """Lista todos os shapes editáveis de um slide."""
        slide = self._get_slide(slide_number)
        if not slide:
            return []
        result = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                result.append({
                    "shape_id": shape.shape_id,
                    "name": shape.name,
                    "is_placeholder": shape.is_placeholder,
                    "placeholder_idx": shape.placeholder_format.idx if shape.is_placeholder else None,
                    "text_preview": shape.text_frame.text[:80].strip(),
                    "width_in": round(shape.width.inches, 2) if shape.width else None,
                    "height_in": round(shape.height.inches, 2) if shape.height else None,
                })
        return result

    # ── Internos ──────────────────────────────────────────────────────────────

    def _get_slide(self, slide_number: int):
        """Retorna o slide pelo número (1-indexed)."""
        idx = slide_number - 1
        if 0 <= idx < len(self.prs.slides):
            return self.prs.slides[idx]
        return None

    def _find_shape_by_name(self, slide, name: str):
        """Encontra shape pelo nome (case-insensitive)."""
        name_lower = name.lower()
        for shape in slide.shapes:
            if shape.name.lower() == name_lower:
                return shape
        # Busca parcial
        for shape in slide.shapes:
            if name_lower in shape.name.lower():
                return shape
        return None

    def _find_body_shape(self, slide):
        """
        Encontra o shape de corpo principal de um slide.
        Prioriza: placeholder idx 1 → maior área com texto → qualquer texto.
        """
        # 1. Placeholder de conteúdo (idx 1)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.idx == 1:
                        return shape
                except Exception:
                    pass

        # 2. Maior shape com texto que NÃO é título
        best = None
        best_area = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.is_placeholder:
                try:
                    if shape.placeholder_format.idx == 0:
                        continue  # pular título
                except Exception:
                    pass
            try:
                area = shape.width.inches * shape.height.inches
                if area > best_area:
                    best_area = area
                    best = shape
            except Exception:
                pass

        return best
