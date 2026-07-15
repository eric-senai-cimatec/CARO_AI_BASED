from pptx import Presentation


def extract_template(template_path: str) -> list:
    prs = Presentation(template_path)
    slides_info = []

    for idx, slide in enumerate(prs.slides, start=1):
        placeholders = []
        title_text = ""
        title_shape = slide.shapes.title

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            is_title = shape == title_shape
            text = shape.text_frame.text.strip()

            info = {
                "type": "title" if is_title else "body",
                "text": text,
                "x": int(shape.left),
                "y": int(shape.top),
                "width": int(shape.width),
                "height": int(shape.height),
            }

            if is_title and text:
                title_text = text

            placeholders.append(info)

        slides_info.append({
            "slide": idx,
            "title": title_text,
            "placeholders": placeholders,
        })

    return slides_info
