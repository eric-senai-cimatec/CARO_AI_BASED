"""
tests/test_e2e.py — Testes de ponta a ponta

Executa o pipeline completo com o FSIPP de exemplo e verifica
que o PPTX é gerado corretamente.

Uso:
    pytest tests/ -v
    pytest tests/test_e2e.py::test_analyst_agent -v
    pytest tests/ -v --tb=short
"""

from __future__ import annotations

import json
import os
import sys
from pathlib import Path

import pytest

# Adicionar root ao path
ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(ROOT))


# ── Fixtures ──────────────────────────────────────────────────────────────────

EXAMPLE_FSIPP = ROOT / "docs" / "exemplo_fsipp.txt"
TEMPLATE = ROOT / "templates" / "CARO_Bobinas.pptx"
OUTPUT_DIR = ROOT / "output" / "tests"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


@pytest.fixture(scope="session")
def fsipp_text() -> str:
    if not EXAMPLE_FSIPP.exists():
        pytest.skip("Arquivo de exemplo não encontrado")
    return EXAMPLE_FSIPP.read_text(encoding="utf-8")


@pytest.fixture(scope="session")
def groq_api_key() -> str:
    key = os.getenv("GROQ_API_KEY", "")
    if not key:
        pytest.skip("GROQ_API_KEY não configurada — skip de testes com LLM")
    return key


@pytest.fixture(scope="session")
def template_model():
    if not TEMPLATE.exists():
        pytest.skip("Template PPTX não encontrado")
    from core.template_reader import TemplateReader
    reader = TemplateReader(cache_dir=ROOT / ".cache")
    return reader.read(TEMPLATE)


# ── Testes de AST ─────────────────────────────────────────────────────────────

class TestAstPptx:
    def test_parse_template(self, template_model):
        ast = template_model.ast
        assert ast.slide_count > 0, "Template deve ter pelo menos 1 slide"
        assert ast.slide_width_in > 0
        assert ast.slide_height_in > 0

    def test_serialize_deserialize(self, template_model):
        import tempfile
        ast = template_model.ast
        with tempfile.NamedTemporaryFile(suffix=".json", delete=False) as f:
            tmp = Path(f.name)
        ast.save_json(tmp)

        from core.ast_pptx import PptxAST
        loaded = PptxAST.load_json(tmp)
        assert loaded.slide_count == ast.slide_count
        assert loaded.file_hash == ast.file_hash
        tmp.unlink()

    def test_get_slide(self, template_model):
        slide = template_model.ast.get_slide(1)
        assert slide is not None
        assert slide.slide_number == 1

    def test_all_slides_have_number(self, template_model):
        for slide in template_model.ast.slides:
            assert slide.slide_number > 0


# ── Testes do Template Reader ─────────────────────────────────────────────────

class TestTemplateReader:
    def test_identity_extracted(self, template_model):
        identity = template_model.identity
        assert identity is not None

    def test_capacities(self, template_model):
        for num, cap in template_model.slide_capacities.items():
            assert cap.slide_number == num
            assert cap.body_max_chars > 0

    def test_editable_fields(self, template_model):
        fields = template_model.editable_fields
        assert len(fields) > 0


# ── Testes do Region Detector ─────────────────────────────────────────────────

class TestRegionDetector:
    def test_detect_slide(self, template_model):
        from core.region_detector import RegionDetector, RegionType
        detector = RegionDetector(
            slide_width_in=template_model.ast.slide_width_in,
            slide_height_in=template_model.ast.slide_height_in,
        )
        slide = template_model.ast.slides[0]
        region_map = detector.detect(slide)
        assert region_map.slide_number == 1
        assert isinstance(region_map.regions, list)


# ── Testes de ferramentas ─────────────────────────────────────────────────────

class TestTextTools:
    def test_truncate_smart(self):
        from tools.text_tools import truncate_smart
        text = "Esta é uma frase longa que deve ser truncada inteligentemente"
        result = truncate_smart(text, 30)
        assert len(result) <= 31  # margem do sufixo
        assert not result.endswith(" ")

    def test_normalize_bullets(self):
        from tools.text_tools import normalize_bullets
        bullets = ["• Primeiro item", "- Segundo item", "3. Terceiro item", ""]
        result = normalize_bullets(bullets)
        assert len(result) == 3
        assert not any(b.startswith("•") for b in result)
        assert not any(b.startswith("-") for b in result)

    def test_sanitize_for_pptx(self):
        from tools.text_tools import sanitize_for_pptx
        text = 'Texto com "aspas tipográficas" e – travessão'
        result = sanitize_for_pptx(text)
        assert '"' in result  # aspas normais
        assert "-" in result  # travessão normalizado

    def test_split_into_chunks(self):
        from tools.text_tools import split_into_chunks
        items = list(range(14))
        chunks = split_into_chunks(items, max_per_chunk=6)
        assert len(chunks) == 3
        assert len(chunks[0]) == 6
        assert len(chunks[2]) == 2


class TestPdfTools:
    def test_read_txt_file(self):
        from tools.pdf_tools import read_fsipp
        text = read_fsipp(EXAMPLE_FSIPP)
        assert len(text) > 100
        assert "TechBob" in text

    def test_read_direct_text(self):
        from tools.pdf_tools import read_fsipp
        direct = "Este é um texto direto de mais de 50 caracteres para teste do framework CARO"
        result = read_fsipp(direct)
        assert result == direct


class TestValidationTools:
    def test_api_key_missing(self):
        from tools.validation_tools import validate_api_key
        ok, msg = validate_api_key("")
        assert not ok
        assert "GROQ_API_KEY" in msg

    def test_api_key_valid(self):
        from tools.validation_tools import validate_api_key
        ok, msg = validate_api_key("gsk_" + "x" * 40)
        assert ok

    def test_template_not_found(self):
        from tools.validation_tools import validate_template
        ok, msg = validate_template("/nao/existe.pptx")
        assert not ok

    def test_template_valid(self):
        if not TEMPLATE.exists():
            pytest.skip("Template não encontrado")
        from tools.validation_tools import validate_template
        ok, msg = validate_template(TEMPLATE)
        assert ok


# ── Testes dos Agentes (requerem GROQ_API_KEY) ────────────────────────────────

class TestAnalystAgent:
    def test_run(self, fsipp_text, groq_api_key):
        from agents.analyst_agent import AnalystAgent
        agent = AnalystAgent("test_analyst")
        result = agent.run(fsipp_text)
        assert result.success, f"AnalystAgent falhou: {result.error}"
        ctx = result.content
        assert ctx.cliente.empresa != ""
        assert ctx.projeto.titulo != ""
        assert len(ctx.projeto.beneficios) > 0

    def test_confidence_score(self, fsipp_text, groq_api_key):
        from agents.analyst_agent import AnalystAgent
        agent = AnalystAgent("test_analyst")
        result = agent.run(fsipp_text)
        assert result.success
        assert result.content.confidence_score > 0.5


class TestPlannerAgent:
    def test_run(self, fsipp_text, groq_api_key, template_model):
        from agents.analyst_agent import AnalystAgent
        from agents.planner_agent import PlannerAgent

        analyst = AnalystAgent("test_analyst")
        ctx = analyst.run(fsipp_text).content

        planner = PlannerAgent("test_planner")
        result = planner.run(ctx, template_model)
        assert result.success
        plan = result.content
        assert len(plan.instructions) > 0
        assert len(plan.get_active_instructions()) > 5


class TestWriterAgent:
    def test_run(self, fsipp_text, groq_api_key, template_model):
        from agents.analyst_agent import AnalystAgent
        from agents.planner_agent import PlannerAgent
        from agents.writer_agent import WriterAgent

        analyst = AnalystAgent("test_analyst")
        ctx = analyst.run(fsipp_text).content

        planner = PlannerAgent("test_planner")
        plan = planner.run(ctx, template_model).content

        writer = WriterAgent("test_writer")
        result = writer.run(ctx, plan, template_model)
        assert result.success
        output = result.content
        assert len(output.slide_contents) > 0

        # Slide 4 (demandante) deve ter key_values
        sc4 = output.get(4)
        if sc4:
            assert sc4.content_type in ("key_value", "bullets", "text", "skip")


# ── Teste E2E completo ────────────────────────────────────────────────────────

class TestEndToEnd:
    def test_full_pipeline(self, groq_api_key):
        """Teste completo: FSIPP TXT → PPTX gerado."""
        if not TEMPLATE.exists():
            pytest.skip("Template não encontrado")

        from agents.orchestrator import Orchestrator

        output_path = OUTPUT_DIR / "test_output.pptx"

        orch = Orchestrator(
            template_path=TEMPLATE,
            groq_api_key=groq_api_key,
            audit_enabled=True,
        )

        result = orch.run(
            input_source=EXAMPLE_FSIPP,
            output_path=output_path,
            skip_preflight=False,
        )

        print(f"\n{result.summary()}")

        assert result.success, f"Pipeline falhou:\n{chr(10).join(result.errors)}"
        assert output_path.exists(), "Arquivo de saída não foi criado"
        assert output_path.stat().st_size > 50_000, "Arquivo gerado muito pequeno"

        # Verificar que o PPTX pode ser aberto
        from pptx import Presentation
        prs = Presentation(str(output_path))
        assert len(prs.slides) > 10


if __name__ == "__main__":
    # Rodar um teste específico manualmente
    import subprocess
    subprocess.run([sys.executable, "-m", "pytest", __file__, "-v", "--tb=short"])
