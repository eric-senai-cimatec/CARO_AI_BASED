from pathlib import Path
from pptx import Presentation
from ppt.registry import get as get_renderer


def render(template_path: str, content: dict) -> None:
    prs = Presentation(template_path)

    for slide_data in content.get("slides", []):
        slide_num = slide_data.get("slide", 0) - 1
        if slide_num < 0 or slide_num >= len(prs.slides):
            continue

        slide = prs.slides[slide_num]

        title_text = slide_data.get("title", "")
        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            _set_text(title_shape.text_frame, title_text)

        layout = slide_data.get("layout", "bullet")
        renderer = get_renderer(layout)
        content_data = slide_data.get("content", {})
        renderer.render(slide, content_data)

    output_dir = Path("output")
    output_dir.mkdir(exist_ok=True)
    output_path = output_dir / "apresentacao_caro.pptx"
    prs.save(str(output_path))


def _set_text(text_frame: object, new_text: str) -> None:
    if not text_frame.paragraphs:
        return
    paragraph = text_frame.paragraphs[0]
    if paragraph.runs:
        paragraph.runs[0].text = new_text
    else:
        paragraph.text = new_text
