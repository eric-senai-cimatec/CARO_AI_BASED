from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Pt

from ppt.renderers.base import BaseRenderer


_COLORS = {
    "primary": RGBColor(0x1B, 0x3A, 0x5C),
    "accent": RGBColor(0x2E, 0x86, 0xAB),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "dark": RGBColor(0x33, 0x33, 0x33),
    "gray": RGBColor(0x88, 0x88, 0x88),
    "line": RGBColor(0x99, 0x99, 0x99),
}


class TableRenderer(BaseRenderer):
    def render(self, slide: Slide, content: dict) -> None:
        if not content:
            return

        prs = slide.part.package.presentation_part.presentation
        sw = prs.slide_width
        sh = prs.slide_height

        headers = content.get("headers", [])
        rows_data = content.get("rows", [])
        if not headers or not rows_data:
            return

        ml = int(sw * 0.05)
        mr = int(sw * 0.05)
        mt = int(sh * 0.18)
        mb = int(sh * 0.06)

        n_rows = 1 + len(rows_data)
        n_cols = len(headers)

        table_width = sw - ml - mr
        table_height = sh - mt - mb
        if table_height <= 0 or table_width <= 0:
            return

        col_widths = self._calc_col_widths(table_width, n_cols)

        table_shape = slide.shapes.add_table(n_rows, n_cols, ml, mt, table_width, table_height)
        table = table_shape.table

        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = cw

        for ci, header in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(header)
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = _COLORS["white"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = _COLORS["primary"]

        for ri, row in enumerate(rows_data):
            for ci in range(n_cols):
                val = str(row[ci]) if ci < len(row) else ""
                cell = table.cell(ri + 1, ci)
                cell.text = ""
                p = cell.text_frame.paragraphs[0]
                p.text = val
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                run = p.runs[0]
                run.font.size = Pt(9)
                run.font.color.rgb = _COLORS["dark"]
                if ri % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF7, 0xFA)

        for ri in range(n_rows):
            for ci in range(n_cols):
                cell = table.cell(ri, ci)
                cell.margin_left = Emu(45720)
                cell.margin_right = Emu(45720)
                cell.margin_top = Emu(18288)
                cell.margin_bottom = Emu(18288)

    def _calc_col_widths(self, table_width, n_cols):
        col_w = table_width // n_cols
        widths = [col_w] * n_cols
        remainder = table_width - col_w * n_cols
        if remainder and widths:
            widths[-1] += remainder
        return widths
