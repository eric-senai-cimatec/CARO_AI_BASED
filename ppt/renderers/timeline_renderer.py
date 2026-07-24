from pptx.slide import Slide
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN

from ppt.renderers.base import BaseRenderer


class TimelineRenderer(BaseRenderer):
    def render(self, slide: Slide, content: dict) -> None:
        prs = slide.part.package.presentation_part.presentation
        sw = prs.slide_width
        sh = prs.slide_height

        ml = int(sw * 0.05)
        mr = int(sw * 0.05)
        mt = int(sh * 0.15)

        milestones = content.get("milestones", []) if content else []
        if not milestones:
            return

        y = mt + int(sh * 0.05)
        spacing = (sw - ml - mr) / (len(milestones) + 1)

        line = slide.shapes.add_connector(1, ml, y, int(sw - mr), y)
        line.line.color.rgb = _c("accent")
        line.line.width = Pt(2)

        for i, m in enumerate(milestones):
            cx = int(ml + spacing * (i + 1))
            label = m if isinstance(m, str) else m.get("label", "")

            txbox = slide.shapes.add_textbox(
                cx - Emu(457200), y + Emu(91440),
                Emu(914400), Emu(137160),
            )
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = label
            p.font.size = Pt(9)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER


def _c(name: str):
    from pptx.dml.color import RGBColor
    palette = {
        "accent": RGBColor(0x2E, 0x86, 0xAB),
        "dark": RGBColor(0x1B, 0x3A, 0x5C),
    }
    return palette.get(name, RGBColor(0, 0, 0))
