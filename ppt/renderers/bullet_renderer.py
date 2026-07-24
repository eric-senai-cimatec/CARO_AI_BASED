from pptx.slide import Slide

from ppt.renderers.base import BaseRenderer


class BulletRenderer(BaseRenderer):
    def render(self, slide: Slide, content: dict) -> None:
        if not content:
            return

        body_text = content.get("body", "")
        if not body_text:
            return

        title_shape = slide.shapes.title
        body_filled = False

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == title_shape:
                continue
            if not body_filled:
                _set_text(shape.text_frame, body_text)
                body_filled = True


def _set_text(text_frame: object, new_text: str) -> None:
    if not text_frame.paragraphs:
        return
    paragraph = text_frame.paragraphs[0]
    if paragraph.runs:
        paragraph.runs[0].text = new_text
    else:
        paragraph.text = new_text
