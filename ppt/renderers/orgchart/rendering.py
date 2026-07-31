from __future__ import annotations

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Pt

from ppt.renderers.orgchart.models import OrgChart, OrgNode
from ppt.renderers.orgchart.utils import (
    ColorPalette,
    Dimensions,
    _DEFAULT_PALETTE,
    get_level_color,
)


class OrgChartRenderer:
    def __init__(
        self,
        slide: Slide,
        dims: Dimensions,
        palette: ColorPalette = _DEFAULT_PALETTE,
    ) -> None:
        self.slide = slide
        self.dims = dims
        self.palette = palette

    def render(self, chart: OrgChart) -> None:
        if not chart.root:
            return
        for node in chart.nodes:
            self._draw_node(node, chart.max_depth)
        for node in chart.nodes:
            for child in node.children:
                self._draw_connector(node, child)

    def _draw_node(self, node: OrgNode, max_depth: int) -> None:
        is_root = node.parent_id is None
        bg_color, text_color = get_level_color(
            node.level, is_root, max_depth, self.palette
        )
        left = int(node.x - node.width / 2.0)
        top = int(node.y)
        width = int(node.width)
        height = int(node.height)

        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = self.palette.border
        shape.line.width = Pt(1.0)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(45720)
        tf.margin_right = Emu(45720)
        tf.margin_top = Emu(18288)
        tf.margin_bottom = Emu(9144)

        p = tf.paragraphs[0]
        p.text = node.text
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(int(self.dims.font_size))
        run.font.color.rgb = text_color
        run.font.bold = is_root

    def _draw_connector(self, parent: OrgNode, child: OrgNode) -> None:
        x1 = int(parent.x)
        y1 = int(parent.y + parent.height)
        x2 = int(child.x)
        y2 = int(child.y)
        mid_y = (y1 + y2) // 2

        conn_left = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, x1, y1, x1, mid_y
        )
        conn_left.line.color.rgb = self.palette.line
        conn_left.line.width = Pt(1.0)

        conn_right = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, x1, mid_y, x2, mid_y
        )
        conn_right.line.color.rgb = self.palette.line
        conn_right.line.width = Pt(1.0)

        conn_down = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, x2, mid_y, x2, y2
        )
        conn_down.line.color.rgb = self.palette.line
        conn_down.line.width = Pt(1.0)
