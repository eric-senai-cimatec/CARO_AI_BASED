import re

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Pt

from ppt.renderers.base import BaseRenderer


_COLORS = {
    "primary": RGBColor(0x1B, 0x3A, 0x5C),
    "accent": RGBColor(0x2E, 0x86, 0xAB),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "dark": RGBColor(0x33, 0x33, 0x33),
    "gray": RGBColor(0x88, 0x88, 0x88),
    "line": RGBColor(0x99, 0x99, 0x99),
    "muted_bg": RGBColor(0xEB, 0xF5, 0xFB),
}


class GanttRenderer(BaseRenderer):
    def render(self, slide: Slide, content: dict) -> None:
        if not content:
            return

        prs = slide.part.package.presentation_part.presentation
        sw = prs.slide_width
        sh = prs.slide_height

        ml = int(sw * 0.04)
        mr = int(sw * 0.04)
        mt = int(sh * 0.16)
        mb = int(sh * 0.06)

        start_label = content.get("start_label", "")
        phases = content.get("phases", [])
        if not phases:
            return

        parsed = []
        for p in phases:
            name = p.get("name", "")
            start_str = str(p.get("start", "") or "")
            end_str = str(p.get("end", "") or "")
            tasks = p.get("tasks", [])
            if isinstance(tasks, str):
                tasks = [tasks]
            start_num = self._parse_month(start_str)
            end_num = self._parse_month(end_str)
            if end_num is None and start_num is not None:
                end_num = start_num
            if start_num is not None and end_num is not None:
                parsed.append((name, start_str, end_str, tasks, start_num, end_num))

        if not parsed:
            return

        min_month = min(s for _, _, _, _, s, _ in parsed)
        max_month = max(e for _, _, _, _, _, e in parsed)
        total_months = max_month - min_month + 1

        label_w = int(sw * 0.22)
        chart_left = ml + label_w + int(sw * 0.02)
        chart_w = sw - chart_left - mr

        header_h = int(sh * 0.04)
        row_h = self._calc_row_height(sh)
        y = mt

        if start_label:
            self._add_textbox(
                slide, ml, y, sw - ml - mr, int(sh * 0.035),
                f"Início: {start_label}",
                font_size=Pt(10), color=_COLORS["dark"], bold=True,
            )
            y += int(sh * 0.035)

        month_w = chart_w // total_months

        for i in range(total_months):
            month_num = min_month + i
            x = chart_left + i * month_w
            hdr = slide.shapes.add_shape(1, x, y, month_w, header_h)
            hdr.fill.solid()
            hdr.fill.fore_color.rgb = _COLORS["primary"]
            hdr.line.fill.background()
            self._add_textbox(
                slide, x, y, month_w, header_h,
                str(month_num),
                font_size=Pt(9), color=_COLORS["white"], bold=True,
                alignment=PP_ALIGN.CENTER,
            )

        sep_y = y + header_h

        for i in range(total_months + 1):
            x = chart_left + i * month_w
            line = slide.shapes.add_connector(1, x, sep_y, x, sep_y + len(parsed) * row_h)
            line.line.color.rgb = _COLORS["line"]
            line.line.width = Pt(0.5)

        sep_line = slide.shapes.add_connector(1, ml, sep_y, ml + label_w + chart_w, sep_y)
        sep_line.line.color.rgb = _COLORS["primary"]
        sep_line.line.width = Pt(1.5)

        for i, (name, start_str, end_str, tasks, s, e) in enumerate(parsed):
            row_top = sep_y + i * row_h
            bar_y = row_top + int(row_h * 0.12)
            bar_h = int(row_h * 0.38) if tasks else int(row_h * 0.55)

            self._add_textbox(
                slide, ml, row_top, label_w, row_h,
                name, font_size=Pt(9), bold=True,
                color=_COLORS["dark"],
            )

            bar_left = chart_left + (s - min_month) * month_w
            bar_width = (e - s + 1) * month_w
            bar = slide.shapes.add_shape(1, bar_left, bar_y, bar_width, bar_h)
            bar.fill.solid()
            bar.fill.fore_color.rgb = _COLORS["accent"]
            bar.line.fill.background()

            if tasks:
                task_top = bar_y + bar_h + Emu(4572)
                task_h = row_h - (task_top - row_top) - Emu(9144)
                if task_h > 0:
                    tasks_text = "\n".join(f"  {t}" for t in tasks if t)
                    self._add_textbox(
                        slide, chart_left, task_top, chart_w, task_h,
                        tasks_text, font_size=Pt(6), color=_COLORS["dark"],
                    )

            if start_str:
                self._add_textbox(
                    slide, bar_left, bar_y + bar_h + Emu(4572),
                    month_w, Emu(68580),
                    start_str, font_size=Pt(6), color=_COLORS["gray"],
                    alignment=PP_ALIGN.CENTER,
                )

            if end_str:
                self._add_textbox(
                    slide, bar_left + bar_width - month_w, bar_y + bar_h + Emu(4572),
                    month_w, Emu(68580),
                    end_str, font_size=Pt(6), color=_COLORS["gray"],
                    alignment=PP_ALIGN.CENTER,
                )

            row_line = slide.shapes.add_connector(
                1, ml, row_top + row_h,
                ml + label_w + chart_w, row_top + row_h,
            )
            row_line.line.color.rgb = _COLORS["line"]
            row_line.line.width = Pt(0.5)

    def _calc_row_height(self, sh: int) -> int:
        return int(sh * 0.065)

    def _parse_month(self, text) -> int | None:
        if not text:
            return None
        m = re.search(r"(\d+)", text)
        return int(m.group(1)) if m else None

    def _add_textbox(
        self, slide, left, top, width, height, text,
        font_size=Pt(10), color=None, bold=False,
        alignment=PP_ALIGN.LEFT,
    ):
        if not text:
            return None
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        run = p.runs[0]
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        p.alignment = alignment
        return txbox