from pathlib import Path

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.util import Emu, Pt

from ppt.renderers.base import BaseRenderer


_COLORS = {
    "primary": RGBColor(0x1B, 0x3A, 0x5C),
    "accent": RGBColor(0x2E, 0x86, 0xAB),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "dark": RGBColor(0x33, 0x33, 0x33),
    "gray": RGBColor(0x88, 0x88, 0x88),
    "muted_bg": RGBColor(0xEB, 0xF5, 0xFB),
    "line": RGBColor(0x99, 0x99, 0x99),
    "milestone_line": RGBColor(0x4A, 0x90, 0xD9),
}


class WorkflowRenderer(BaseRenderer):
    def __init__(self, icon_dir: str | None = None) -> None:
        self._icon_dir = Path(icon_dir) if icon_dir else (
            Path(__file__).parent.parent.parent / "icons"
        )

    def render(self, slide: Slide, content: dict) -> None:
        if not content:
            return

        prs = slide.part.package.presentation_part.presentation
        sw = prs.slide_width
        sh = prs.slide_height

        ml = int(sw * 0.04)
        mr = int(sw * 0.04)
        usable_w = sw - ml - mr

        mt = int(sh * 0.18)
        mb = int(sh * 0.06)

        timeline_data = content.get("timeline", {})
        steps = content.get("steps", [])

        time_h = int(sh * 0.18)
        self._draw_timeline(slide, timeline_data, ml, mt, usable_w, time_h)

        time_bottom = mt + time_h
        steps_top = time_bottom + int(sh * 0.04)
        steps_area_h = sh - steps_top - mb

        if steps and steps_area_h > 0:
            n = len(steps)
            gap = int(usable_w * 0.02)
            box_w = (usable_w - gap * (n - 1)) // n
            box_h = steps_area_h

            for i, step in enumerate(steps):
                left = ml + i * (box_w + gap)
                self._draw_step_box(slide, step, left, steps_top, box_w, box_h)

            for i in range(n - 1):
                a_left = ml + (i + 1) * (box_w + gap) - gap // 2
                a_y = steps_top + box_h // 2
                self._draw_arrow(slide, a_left, a_y, a_left + gap, a_y)

    def _draw_timeline(self, slide, timeline, left, top, width, height):
        start_label = (timeline.get("start") or "").strip()
        milestones = timeline.get("milestones") or []

        line_y = top + height // 3
        diamond_sz = int(height * 0.25)

        line = slide.shapes.add_connector(1, left, line_y, left + width, line_y)
        line.line.color.rgb = _COLORS["milestone_line"]
        line.line.width = Pt(2.5)

        self._add_textbox(
            slide, left, top, int(width * 0.18), line_y - top,
            start_label,
            font_size=Pt(8), color=_COLORS["gray"],
        )

        if not milestones:
            return

        spacing = width / (len(milestones) + 1)

        for i, ms in enumerate(milestones):
            cx = int(left + spacing * (i + 1))

            diamond = slide.shapes.add_shape(
                5, cx - diamond_sz // 2, line_y - diamond_sz // 2,
                diamond_sz, diamond_sz,
            )
            diamond.fill.solid()
            diamond.fill.fore_color.rgb = _COLORS["accent"]
            diamond.line.fill.background()

            lbl = ms if isinstance(ms, str) else ""
            self._add_textbox(
                slide,
                cx - int(width * 0.08), line_y + diamond_sz // 2 + Emu(27432),
                int(width * 0.16), Emu(137160),
                lbl,
                font_size=Pt(8), color=_COLORS["dark"], bold=True,
                alignment=PP_ALIGN.CENTER,
            )

            vert = slide.shapes.add_connector(
                1, cx, line_y + diamond_sz // 2,
                cx, top + height - Emu(27432),
            )
            vert.line.color.rgb = _COLORS["line"]
            vert.line.width = Pt(1)
            vert.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    def _draw_step_box(self, slide, step, left, top, width, height):
        pad = Emu(45720)

        box = slide.shapes.add_shape(1, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = _COLORS["white"]
        box.line.color.rgb = _COLORS["accent"]
        box.line.width = Pt(1.5)
        box.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        title_text = (step.get("title") or "").strip()
        title_h = int(height * 0.18)

        if title_text:
            header = slide.shapes.add_shape(1, left, top, width, title_h)
            header.fill.solid()
            header.fill.fore_color.rgb = _COLORS["primary"]
            header.line.fill.background()

            self._add_textbox(
                slide,
                left + pad, top + int(title_h * 0.15),
                width - 2 * pad, int(title_h * 0.70),
                title_text,
                font_size=Pt(10), color=_COLORS["white"], bold=True,
                alignment=PP_ALIGN.CENTER,
            )

        icon_name = (step.get("icon") or "").strip()
        icon_top = top + title_h + pad
        icon_h = int(height * 0.20)

        if icon_name:
            icon_path = self._resolve_icon(icon_name)
            if icon_path:
                try:
                    slide.shapes.add_picture(
                        str(icon_path),
                        left + (width - icon_h) // 2,
                        icon_top, icon_h, icon_h,
                    )
                    icon_top += icon_h + pad
                except Exception:
                    pass

        items = step.get("items") or []
        if items:
            i_top = icon_top
            i_h = top + height - i_top - pad
            if i_h > 0:
                self._add_bullet_list(
                    slide, left + pad, i_top, width - 2 * pad, i_h, items,
                )

    def _add_textbox(self, slide, left, top, width, height, text,
                     font_size=Pt(10), color=None, bold=False,
                     alignment=PP_ALIGN.LEFT):
        if not text:
            return None
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        run = p.runs[0]
        run.font.size = font_size
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        p.alignment = alignment
        return txbox

    def _add_bullet_list(self, slide, left, top, width, height, items):
        txbox = slide.shapes.add_textbox(left, top, width, height)
        tf = txbox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(9)
            p.font.color.rgb = _COLORS["dark"]
            p.space_after = Pt(2)
            pPr = p._p.get_or_add_pPr()
            buChar = etree.SubElement(pPr, qn("a:buChar"))
            buChar.set("char", "\u2022")

    def _draw_arrow(self, slide, x1, y1, x2, y2):
        conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
        conn.line.color.rgb = _COLORS["accent"]
        conn.line.width = Pt(1.5)
        tail = etree.SubElement(conn.line._ln, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")

    def _resolve_icon(self, name):
        candidates = [
            self._icon_dir / f"{name}.png",
            self._icon_dir / f"{name}.PNG",
        ]
        for c in candidates:
            if c.exists():
                return c
        return None
