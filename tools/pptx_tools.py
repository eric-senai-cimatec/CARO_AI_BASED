"""
tools/pptx_tools.py — Tools PPTX para os Agentes

Conjunto de funções que os agentes podem invocar via Tool Calling:
  - list_slides: lista slides com título e capacidade
  - read_slide: lê conteúdo atual de um slide
  - write_slide_title: escreve título em um slide
  - write_slide_body: escreve corpo em um slide
  - list_shapes: lista shapes de um slide

Também usado diretamente pelo Orchestrator para construção do PPTX final.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation

from core.shape_editor import ShapeEditor
from core.template_reader import TemplateModel

logger = logging.getLogger(__name__)


class PptxTools:
    """
    Ferramentas de alto nível para manipulação de PPTX.

    Encapsula o ShapeEditor e o TemplateModel para oferecer
    uma interface simples e segura para o Orchestrator.

    Uso:
        tools = PptxTools(prs, template_model)
        tools.write_slide_title(4, "Dados do Demandante")
        tools.write_slide_body(4, bullets=["Empresa: CARO", "Contato: João"])
        tools.save("output/proposta.pptx")
    """

    def __init__(self, presentation: Presentation, template_model: Optional[TemplateModel] = None):
        self.prs = presentation
        self.editor = ShapeEditor(presentation)
        self.model = template_model

    # ── Leitura ───────────────────────────────────────────────────────────────

    def list_slides(self) -> List[Dict[str, Any]]:
        """Lista todos os slides com título e número de shapes."""
        result = []
        for i, slide in enumerate(self.prs.slides, 1):
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.is_placeholder:
                    try:
                        if shape.placeholder_format.idx == 0:
                            title = shape.text_frame.text.strip()[:80]
                            break
                    except Exception:
                        pass
            if not title:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            title = t[:80]
                            break

            cap = self.model.get_capacity(i) if self.model else None
            result.append({
                "slide_number": i,
                "title": title,
                "shape_count": len(slide.shapes),
                "layout": slide.slide_layout.name,
                "body_max_chars": cap.body_max_chars if cap else None,
            })
        return result

    def read_slide(self, slide_number: int) -> Dict[str, Any]:
        """Lê o conteúdo atual de um slide."""
        idx = slide_number - 1
        if not (0 <= idx < len(self.prs.slides)):
            return {"error": f"Slide {slide_number} não existe"}
        slide = self.prs.slides[idx]
        shapes_data = []
        for shape in slide.shapes:
            entry = {"name": shape.name, "type": str(shape.shape_type)}
            if shape.has_text_frame:
                entry["text"] = shape.text_frame.text.strip()[:200]
            shapes_data.append(entry)
        return {"slide_number": slide_number, "shapes": shapes_data}

    def list_shapes(self, slide_number: int) -> List[Dict[str, Any]]:
        """Lista shapes editáveis de um slide."""
        return self.editor.list_editable_shapes(slide_number)

    # ── Escrita ───────────────────────────────────────────────────────────────

    def write_slide_title(self, slide_number: int, text: str) -> bool:
        """Escreve o título de um slide."""
        ok = self.editor.set_slide_title(slide_number, text)
        if ok:
            logger.debug(f"✏️  Slide {slide_number} título: {text[:50]!r}")
        return ok

    def write_slide_body(
        self,
        slide_number: int,
        text: Optional[str] = None,
        bullets: Optional[List[str]] = None,
        key_values: Optional[List[Tuple[str, str]]] = None,
        shape_name: Optional[str] = None,
    ) -> bool:
        """Escreve o corpo de um slide."""
        ok = self.editor.set_slide_body(
            slide_number, text=text,
            bullets=bullets, key_values=key_values,
            shape_name=shape_name,
        )
        if ok:
            preview = (text or " | ".join(bullets or []) or str(key_values))[:50]
            logger.debug(f"✏️  Slide {slide_number} corpo: {preview!r}")
        return ok

    def write_slide_date(self, slide_number: int, date_text: str) -> bool:
        """Escreve o campo de data de um slide."""
        return self.editor.set_date(slide_number, date_text)

    def write_table(self, slide_number: int, data: List[List[str]]) -> bool:
        """Preenche a tabela de um slide."""
        return self.editor.edit_table(slide_number, data)

    def write_shape(self, slide_number: int, shape_name: str, text: str) -> bool:
        """Escreve em um shape específico pelo nome."""
        return self.editor.set_shape_text(slide_number, shape_name, text)

    # ── Persistência ──────────────────────────────────────────────────────────

    def save(self, output_path: str | Path) -> Path:
        """Salva o PPTX editado."""
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        size_kb = path.stat().st_size // 1024
        logger.info(f"💾 Salvo: {path} ({size_kb} KB)")
        return path

    # ── Tool definitions (para agent tool-calling) ────────────────────────────

    def get_tool_definitions(self) -> List[Dict[str, Any]]:
        """Retorna definições de tools no formato Groq/OpenAI."""
        return [
            {
                "type": "function",
                "function": {
                    "name": "list_slides",
                    "description": "Lista todos os slides com título e capacidade",
                    "parameters": {"type": "object", "properties": {}, "required": []},
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "read_slide",
                    "description": "Lê o conteúdo atual de um slide específico",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "slide_number": {"type": "integer", "description": "Número do slide (1-indexed)"}
                        },
                        "required": ["slide_number"],
                    },
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "write_slide_title",
                    "description": "Escreve o título de um slide preservando formatação",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "slide_number": {"type": "integer"},
                            "text": {"type": "string", "description": "Título do slide"},
                        },
                        "required": ["slide_number", "text"],
                    },
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "write_slide_body",
                    "description": "Escreve o corpo de um slide (bullets ou texto)",
                    "parameters": {
                        "type": "object",
                        "properties": {
                            "slide_number": {"type": "integer"},
                            "bullets": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "Lista de bullets (use bullets OU text, não ambos)",
                            },
                            "text": {
                                "type": "string",
                                "description": "Texto corrido (use text OU bullets, não ambos)",
                            },
                        },
                        "required": ["slide_number"],
                    },
                },
            },
        ]
